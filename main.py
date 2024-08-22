import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from PIL import Image, ImageTk, ImageDraw, ImageFont
import io

# Global zoom level
zoom_level = 1.0

# Function to convert slide to image
def slide_to_image(slide):
    slide_img = Image.new('RGB', (960, 540), color=(255, 255, 255))
    draw = ImageDraw.Draw(slide_img)

    for shape in slide.shapes:
        if not shape.has_text_frame and not hasattr(shape, 'image'):
            continue

        if shape.has_text_frame:
            text = shape.text
            left, top, width, height = (shape.left, shape.top, shape.width, shape.height)

            try:
                if shape.text_frame.paragraphs[0].font is not None:
                    pptx_font = shape.text_frame.paragraphs[0].font
                    font_name = pptx_font.name if pptx_font.name else "Arial"
                    font_size = pptx_font.size.pt if pptx_font.size else 12
                    font = ImageFont.truetype(font_name, int(font_size))
                else:
                    font = ImageFont.load_default()
            except OSError:
                font = ImageFont.load_default()

            draw.text((left / 12700, top / 12700), text, fill=(0, 0, 0), font=font)

        if hasattr(shape, 'image'):
            image_stream = io.BytesIO(shape.image.blob)
            img = Image.open(image_stream)
            img.thumbnail((width / 12700, height / 12700))
            slide_img.paste(img, (int(left / 12700), int(top / 12700)))

    return slide_img

# Function to display the slide in Tkinter with zoom
def display_slide(slide, slide_number, total_slides):
    global zoom_level
    img = slide_to_image(slide)

    # Apply zoom to the image
    width, height = img.size
    img = img.resize((int(width * zoom_level), int(height * zoom_level)), Image.LANCZOS)

    img_tk = ImageTk.PhotoImage(img)
    slide_label.config(image=img_tk)
    slide_label.image = img_tk

    # Update the page designator
    page_designator.config(text=f"Slide {slide_number + 1}/{total_slides}")
    # Update the zoom level entry box
    zoom_entry.delete(0, tk.END)
    zoom_entry.insert(0, f"{int(zoom_level * 100)}%")

# Function to increase the zoom level
def zoom_in():
    global zoom_level
    zoom_level += 0.05  # Increase zoom level by 5%
    display_slide(slides[current_slide], current_slide, len(slides))

# Function to decrease the zoom level
def zoom_out():
    global zoom_level
    if zoom_level > 0.05:  # Prevent zoom level from going below 5%
        zoom_level -= 0.05
    display_slide(slides[current_slide], current_slide, len(slides))

# Function to set the zoom level manually
def set_zoom_level(event=None):
    global zoom_level
    try:
        # Get the zoom level from the entry box and convert to a float
        zoom_percent = float(zoom_entry.get().strip('%'))
        zoom_level = zoom_percent / 100.0
        display_slide(slides[current_slide], current_slide, len(slides))
    except ValueError:
        # If the entry is invalid, reset the entry box to the current zoom level
        zoom_entry.delete(0, tk.END)
        zoom_entry.insert(0, f"{int(zoom_level * 100)}%")

# Load PowerPoint file and initialize slides
def load_ppt():
    ppt_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if ppt_path:
        prs = Presentation(ppt_path)
        slides.clear()
        slides.extend(prs.slides)
        display_slide(slides[0], 0, len(slides))

# Navigate to the next slide
def next_slide(event=None):
    global current_slide
    current_slide += 1
    if current_slide >= len(slides):
        current_slide = len(slides) - 1
    display_slide(slides[current_slide], current_slide, len(slides))

# Navigate to the previous slide
def previous_slide(event=None):
    global current_slide
    current_slide -= 1
    if current_slide < 0:
        current_slide = 0
    display_slide(slides[current_slide], current_slide, len(slides))

# Handle mouse wheel scrolling
def mouse_wheel(event):
    if event.delta < 0:
        next_slide()
    elif event.delta > 0:
        previous_slide()

# Initialize main application window
root = tk.Tk()
root.title("PowerPoint Viewer")

# Create a label to hold the slide image
slide_label = tk.Label(root)
slide_label.pack()

# Bind click events to the slide label for navigation
slide_label.bind("<Button-1>", next_slide)  # Left-click to advance to the next slide
slide_label.bind("<Control-Button-1>", previous_slide)  # Ctrl + Left-click to go to the previous slide

# Bind mouse wheel scroll events
root.bind("<MouseWheel>", mouse_wheel)  # For Windows and MacOS
root.bind("<Button-4>", lambda event: previous_slide())  # For Linux (scroll up)
root.bind("<Button-5>", lambda event: next_slide())  # For Linux (scroll down)

# Create buttons for navigation and zoom
prev_button = tk.Button(root, text="Previous", command=previous_slide)
prev_button.pack(side=tk.LEFT, padx=10, pady=10)

next_button = tk.Button(root, text="Next", command=next_slide)
next_button.pack(side=tk.RIGHT, padx=10, pady=10)

zoom_in_button = tk.Button(root, text="Zoom In", command=zoom_in)
zoom_in_button.pack(side=tk.LEFT, padx=10, pady=10)

zoom_out_button = tk.Button(root, text="Zoom Out", command=zoom_out)
zoom_out_button.pack(side=tk.LEFT, padx=10, pady=10)

# Zoom level entry box
zoom_entry = tk.Entry(root, width=5)
zoom_entry.pack(side=tk.LEFT, padx=10, pady=10)
zoom_entry.bind("<Return>", set_zoom_level)

# Page designator label
page_designator = tk.Label(root, text="Slide 0/0")
page_designator.pack()

# Create a menu for loading PowerPoint files
menu = tk.Menu(root)
root.config(menu=menu)
file_menu = tk.Menu(menu, tearoff=0)
menu.add_cascade(label="File", menu=file_menu)
file_menu.add_command(label="Open", command=load_ppt)

# Initialize slide list and current slide index
slides = []
current_slide = 0

# Run the application
root.mainloop()
